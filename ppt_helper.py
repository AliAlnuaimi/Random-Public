import re
from pathlib import Path
from pptx import Presentation
import zipfile
import os
from typing import List, Dict, Optional
import tempfile
import shutil
import time
# Below imports may cause conflicts
import pythoncom
import win32com.client
from .excel_openpyxl_helper import replace_text_in_excel_bytes

def replace_text_in_slides_and_embedded_sheets(pptx_path, replacements):
    # Replace text on main ppt body
    success, count = replace_text(pptx_path, replacements, replace_all=True, match_case=False, whole_word=True)
    print(f"Success: {success}, Replacements: {count}")

    # Extract embedded sheets
    embedded_files = extract_embedded_excel_files(pptx_path)
    updated_files = {}

    for path, content in embedded_files.items():
        success, count, updated_bytes = replace_text_in_excel_bytes(content, replacements)
        if success and updated_bytes:
            updated_files[path] = updated_bytes
            print(f"{path}: {count} replacements made.")

    # Inject modified Excel sheets back into the PowerPoint
    if updated_files:
        replace_embedded_excel_files(pptx_path, updated_files)
        print("Updated embedded sheets successfully.")

        refresh_all_charts_in_pptx(pptx_path)
    else:
        print("No changes were made.")

# To add: exclude specific sheets
# To add: return failure reason in return tuple
def replace_text_old(file_path, replace_dict, replace_all=True, match_case=True, whole_word=False):
    r"""
    Reads a PowerPoint file, replaces text based on replace_dict, and writes back to the file.
    
    Args:
        file_path (str): Path to the PowerPoint file (.pptx or .ppt).
        replace_dict (dict): Dictionary with old_text: new_text pairs for replacement.
        replace_all (bool): If True, replace all occurrences; if False, replace only the first occurrence in the file.
        match_case (bool): If True, case-sensitive; if False, case-insensitive.
        whole_word (bool): If True, replace only whole words; if False, replace all matches.
    
    Returns:
        tuple: (success: bool, total_replacements: int)
    
    Example:
        replace_dict = {"replace_me": "replaced_1", "replace_me_2": "replaced_2"}
        success, count = replace_text_in_ppt(r"G:\Main\Business\Software Company\Products\Batch Document Keyword Replacer\Testing\test.pptx", replace_dict, replace_all=True, match_case=False, whole_word=False)
        print(f"Success: {success}, Replacements: {count}")
    """
    try:
        # Validate file path and ensure it exists
        file_path = Path(file_path)
        if not file_path.exists() or not file_path.suffix.lower() in ('.pptx', '.ppt'):
            return (False, 0)

        # Load the PowerPoint presentation
        presentation = Presentation(file_path)
        total_replacements = 0

        # Iterate through all slides
        for slide in presentation.slides:
            # Iterate through all shapes in the slide
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue

                        original_text = run.text
                        new_text = original_text
                        run_replacements = 0

                        # Process each replacement in the dictionary
                        for old_text, new_text_value in replace_dict.items():
                            if not isinstance(old_text, str) or not isinstance(new_text_value, str):
                                continue

                            # Set up regex flags and pattern
                            flags = 0 if match_case else re.IGNORECASE
                            pattern_text = old_text
                            if whole_word:
                                pattern_text = r'\b' + re.escape(old_text) + r'\b'

                            # Perform replacement
                            if replace_all:
                                new_text, count = re.subn(pattern_text, new_text_value, new_text, flags=flags)
                            else:
                                # Replace only the first occurrence
                                match = re.search(pattern_text, new_text, flags=flags)
                                if match and run_replacements == 0:
                                    new_text = (new_text[:match.start()] + new_text_value + 
                                                new_text[match.end():])
                                    count = 1
                                else:
                                    count = 0

                            run_replacements += count
                            total_replacements += count

                            # Stop processing this run if a replacement was made and replace_all is False
                            if not replace_all and run_replacements > 0:
                                break

                        # Update run text if changed
                        if new_text != original_text:
                            run.text = new_text

                        # Stop processing further runs if a replacement was made and replace_all is False
                        if not replace_all and total_replacements > 0:
                            # Save the presentation and return
                            presentation.save(file_path)
                            return (True, total_replacements)

        # Save the presentation
        presentation.save(file_path)
        return (True, total_replacements)

    except Exception as e:
        print(f"Error processing file: {e}")
        return (False, 0)

def replace_text(file_path, replace_dict, replace_all=True, match_case=True, whole_word=False):
    r"""
    Replaces text in a PowerPoint file, including placeholders like {{replace}}, <replace>, or [replace],
    in slide body and chart titles, removing delimiters and preserving original text formatting.
    
    Args:
        file_path: Path to the PowerPoint file.
        replace_dict: Dictionary mapping old text (e.g., 'replace' for '{{replace}}') to new text.
        replace_all: If True, replace all occurrences; if False, replace only the first.
        match_case: If True, match case; if False, ignore case.
        whole_word: If True, match whole words only for non-placeholder text.
    
    Returns:
        tuple: (success: bool, total_replacements: int)
    """
    try:
        file_path = Path(file_path)
        if not file_path.exists() or file_path.suffix.lower() not in ('.pptx', '.ppt'):
            print(f"File {file_path} does not exist or is not a valid PowerPoint file.")
            return (False, 0)
        
        if not replace_dict:
            print("No keys provided to replace")
            return (False, 0)

        presentation = Presentation(file_path)
        total_replacements = 0

        # Patterns for different placeholder formats
        placeholder_patterns = [
            (r'\{\{([^}]+)\}\}', lambda x: '{{' + x + '}}'),  # {{text}}
            (r'\<([^>]+)\>', lambda x: '<' + x + '>'),        # <text>
            (r'\[([^\]]+)\]', lambda x: '[' + x + ']')        # [text]
        ]

        def replace_in_text(full_text, runs, replace_dict, replace_all, match_case, whole_word):
            """Helper function to process text in runs while preserving formatting."""
            updated = False
            total_count = 0

            for old_text, new_text_value in replace_dict.items():
                if not isinstance(old_text, str) or not isinstance(new_text_value, str):
                    print(f"Skipping invalid replace pair: {old_text} -> {new_text_value}")
                    continue

                old_text = old_text.strip()
                if not old_text:
                    print("Skipping empty old_text")
                    continue

                # Initialize count for this replacement
                count = 0

                # Try each placeholder pattern and plain text
                patterns = [(re.escape(old_text), old_text)]  # Plain text pattern
                for pattern, format_text in placeholder_patterns:
                    patterns.append((pattern.replace(r'[^}]+', re.escape(old_text)), format_text(old_text)))

                flags = 0 if match_case else re.IGNORECASE

                for pattern, original_text in patterns:
                    # Apply word boundaries only for plain text (not placeholders)
                    if whole_word and original_text == old_text and not any(
                        original_text.startswith(c) or original_text.endswith(c) for c in ('<', '{', '[', '>', '}', ']')
                    ):
                        pattern = r'\b' + pattern + r'\b'

                    # Debug: Log patterns and text
                    #print(f"Searching for pattern '{pattern}' (original: '{original_text}') in text '{full_text}'")

                    # Check for matches
                    matches = list(re.finditer(pattern, full_text, flags=flags))
                    if not matches:
                        continue

                    if replace_all:
                        count = len(matches)
                    else:
                        matches = matches[:1]
                        count = 1 if matches else 0

                    if count == 0:
                        continue

                    total_count += count
                    updated = True

                    # Process each run to preserve formatting
                    new_runs_text = []
                    current_pos = 0

                    for run in runs:
                        run_text = run.text
                        run_start = full_text.find(run_text, current_pos)
                        if run_start == -1:
                            new_runs_text.append((run_text, run.font))
                            continue
                        run_end = run_start + len(run_text)
                        new_run_text = ""

                        for match in matches:
                            match_start, match_end = match.start(), match.end()
                            if match_start >= run_end or match_end <= run_start:
                                continue
                            # Adjust match positions relative to run
                            rel_start = max(0, match_start - run_start)
                            rel_end = min(len(run_text), match_end - run_start)
                            new_run_text += run_text[:rel_start] + new_text_value + run_text[rel_end:]

                        if not new_run_text:
                            new_run_text = run_text
                        new_runs_text.append((new_run_text, run.font))
                        current_pos = run_end

                    # Update runs
                    for i, run in enumerate(runs):
                        if i < len(new_runs_text):
                            run.text = new_runs_text[i][0]
                        else:
                            run.text = ""

                    # Remove empty runs
                    for run in runs[::-1]:
                        if not run.text:
                            run._r.getparent().remove(run._r)

                    if not replace_all:
                        break

            return updated, total_count

        for slide_idx, slide in enumerate(presentation.slides):
            print(slide_idx)
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle text in regular shapes (text frames)
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        if not paragraph.runs:
                            print(f"Slide {slide_idx}, Shape {shape_idx}, Para {para_idx}: Empty paragraph, skipping.")
                            continue

                        full_text = "".join(run.text for run in paragraph.runs)
                        updated, count = replace_in_text(full_text, paragraph.runs, replace_dict, replace_all, match_case, whole_word)
                        total_replacements += count
                        if updated and not replace_all and total_replacements > 0:
                            presentation.save(file_path)
                            print(f"Saved file after {total_replacements} replacements")
                            return (True, total_replacements)

                # Handle text in chart titles
                if shape.has_chart:
                    chart = shape.chart
                    if chart.has_title and chart.chart_title.has_text_frame:
                        title_text_frame = chart.chart_title.text_frame
                        for para_idx, paragraph in enumerate(title_text_frame.paragraphs):
                            if not paragraph.runs:
                                print(f"Slide {slide_idx}, Shape {shape_idx}, Chart Title, Para {para_idx}: Empty paragraph, skipping.")
                                continue

                            full_text = "".join(run.text for run in paragraph.runs)
                            updated, count = replace_in_text(full_text, paragraph.runs, replace_dict, replace_all, match_case, whole_word)
                            total_replacements += count
                            if updated and not replace_all and total_replacements > 0:
                                presentation.save(file_path)
                                print(f"Saved file after {total_replacements} replacements")
                                return (True, total_replacements)

        presentation.save(file_path)
        print(f"Saved file with {total_replacements} total replacements")
        return (True, total_replacements)

    except Exception as e:
        print(f"Error processing file: {e}")
        return (False, 0)
    
def get_embedded_excel_paths(pptx_path: str) -> List[str]:
    """
    Returns a list of embedded Excel file paths inside the .pptx.
    Example: ['ppt/embeddings/oleObject1.xlsx']
    """
    with zipfile.ZipFile(pptx_path, 'r') as zip_file:
        return [f for f in zip_file.namelist() if f.startswith("ppt/embeddings/") and f.endswith(".xlsx")]

def extract_embedded_excel_files(pptx_path: str) -> Dict[str, bytes]:
    """
    Returns a dict of embedded Excel file paths mapped to their raw bytes.
    """
    with zipfile.ZipFile(pptx_path, 'r') as zip_file:
        return {
            f: zip_file.read(f)
            for f in zip_file.namelist()
            if f.startswith("ppt/embeddings/") and f.endswith(".xlsx")
        }

def replace_embedded_excel_files(pptx_path: str, updated_files: dict):
    """
    Safely replaces embedded Excel files in a PPTX by rewriting the archive.
    Preserves structure and avoids corruption.
    """
    temp_fd, temp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(temp_fd)

    try:
        with zipfile.ZipFile(pptx_path, 'r') as src_zip:
            with zipfile.ZipFile(temp_path, 'w', compression=zipfile.ZIP_DEFLATED) as dst_zip:
                for item in src_zip.infolist():
                    if item.filename in updated_files:
                        # Replace with updated content
                        dst_zip.writestr(item.filename, updated_files[item.filename])
                    else:
                        # Copy original content
                        data = src_zip.read(item.filename)
                        dst_zip.writestr(item, data)

        # Replace original PPTX
        shutil.move(temp_path, pptx_path)
        print("✔ Embedded Excel files replaced successfully.")

    except Exception as e:
        print(f"❌ Error rebuilding PPTX: {e}")
        if os.path.exists(temp_path):
            os.remove(temp_path)

def refresh_all_charts_in_pptx(pptx_path: str, save: bool = True) -> bool:
    """
    Opens a PowerPoint file and refreshes all charts by accessing their data links.

    Args:
        pptx_path (str): Full path to the PowerPoint file.
        save (bool): Whether to save the file after refreshing.

    Returns:
        bool: True if success, False if error.
    """
    if not os.path.isfile(pptx_path):
        print(f"File not found: {pptx_path}")
        return False

    try:
        # Start PowerPoint COM
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        # Open the presentation
        presentation = ppt_app.Presentations.Open(pptx_path, WithWindow=False)

        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasChart:
                    try:
                        # Access chart data and force it to load/update
                        shape.Chart.ChartData.Activate()
                        shape.Chart.Refresh()
                        time.sleep(0.1)
                    except Exception as chart_err:
                        print(f"Failed to refresh chart on slide {slide.SlideIndex}: {chart_err}")

        if save:
            presentation.Save()

        presentation.Close()
        ppt_app.Quit()
        pythoncom.CoUninitialize()

        print("✅ All charts refreshed successfully.")
        return True

    except Exception as e:
        print(f"❌ Error refreshing charts: {e}")
        return False